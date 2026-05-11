#!/usr/bin/env python3
"""
Generate FOUNTAIN Pitch Deck - Notion Structure Version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors - Brand dark blue
BG_DARK = RGBColor(0, 5, 16)  # #000510
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(255, 255, 255)  # Use white for text on dark bg
GRAY = RGBColor(255, 255, 255)
GRAY2 = RGBColor(255, 255, 255)
GRAY3 = RGBColor(255, 255, 255)

def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Helvetica Neue'
    p.alignment = align
    return shape

def set_slide_bg(slide):
    """Set slide background to dark blue gradient"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 16, 48)  # #001030

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    # ═══════════ SLIDE 1: COVER ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.4, "🔹 FOUNTAIN", 12, BLACK, True)
    add_text(s, 0.9, 2.2, 11, 1.2, "Physical AI\nBuilder Accelerator", 64, BLACK, True)
    add_text(s, 0.9, 4.2, 10, 0.5, "Bay Area × Shenzhen", 28, GRAY)
    add_text(s, 0.9, 5.0, 10, 0.5, "FOUNTAIN brings AI into the physical world, with builders.", 18, GRAY2)
    
    # ═══════════ SLIDE 2: WHY PHYSICAL AI ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "WHY PHYSICAL AI", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.8, "Physical AI is the next distribution\nlayer of intelligence.", 40, BLACK, True)
    
    points = [
        "AI has scaled rapidly in the digital world. The physical world is just beginning.",
        "AI will move through devices, environments and systems, becoming part of everyday life.",
        "The real world provides data, context, and feedback for models and agents to evolve."
    ]
    for i, pt in enumerate(points):
        add_text(s, 0.9, 2.6 + i * 0.7, 11, 0.6, f"→  {pt}", 17, GRAY)
    
    add_text(s, 0.9, 5.0, 11, 0.3, "Physical AI includes:", 11, GRAY3, True)
    add_text(s, 0.9, 5.35, 11, 0.4, "AI devices, wearables, sensors, embodied systems and robotics.", 15, GRAY)
    
    # ═══════════ SLIDE 3: THE GAP ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "THE REAL-WORLD GAP", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.8, "The gap between demo and\nreal-world product is still large.", 40, BLACK, True)
    
    gaps = [
        "AI-native builders will create the best Physical AI products.\nBut building models or apps ≠ shipping physical systems.",
        "Moving from prototype to mass production requires:\nsystem integration, supply chain, manufacturing and iteration.",
        "US-based teams often need Shenzhen & Asia\nfor real supply chain and production.",
        "Great models are not enough.\nReal-world systems must be built and shipped."
    ]
    for i, g in enumerate(gaps):
        row, col = i // 2, i % 2
        add_text(s, 0.9 + col * 6, 3.0 + row * 1.5, 5.5, 1.2, g, 15, GRAY)
    
    # ═══════════ SLIDE 4: THE SYSTEM ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 5, 0.3, "THE FOUNTAIN SYSTEM", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "A Physical AI product accelerator", 40, BLACK, True)
    add_text(s, 0.9, 1.85, 11, 0.4, "Built between Silicon Valley and Shenzhen.", 17, GRAY2)
    
    add_text(s, 0.9, 2.8, 5.5, 0.4, "🔶 Accelerator Engine", 18, BLACK, True)
    add_text(s, 0.9, 3.25, 5.5, 0.7, "Scaling AI-native teams from prototype to real-world product.\nAI demo → Physical AI product", 15, GRAY)
    
    add_text(s, 7, 2.8, 5.5, 0.4, "🔶 Venture Studio", 18, BLACK, True)
    add_text(s, 7, 3.25, 5.5, 0.7, "Creating new Physical AI companies from the ground up.\nIdea → Company", 15, GRAY)
    
    add_text(s, 0.9, 4.8, 11, 0.5, "\"We don't just advise. We build alongside founders.\"", 18, GRAY2, align=PP_ALIGN.CENTER)
    
    # ═══════════ SLIDE 5: PROGRAM OVERVIEW ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 5, 0.3, "HOW THE PROGRAM WORKS", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "Builders typically work with us for 6 months.", 40, BLACK, True)
    
    phases = [
        ("Phase 1", "San Francisco", "PMF & Product Alignment", "Model capabilities, user value, MP feasibility"),
        ("Phase 2", "Shenzhen", "Shenzhen Immersion", "Factories, labs, supply chain networks"),
        ("Phase 3", "SZ & SF", "Hardware Development", "ID, electronics, PCBA, embedded systems"),
        ("Phase 4", "Shenzhen", "EVT Prototypes", "3-5 EVT units. Tested. Logged. Reviewed."),
        ("Phase 5", "San Francisco", "GTM & Launch", "User testing, storytelling, go-to-market")
    ]
    for i, (num, loc, title, desc) in enumerate(phases):
        left = 0.5 + i * 2.5
        add_text(s, left, 2.5, 2.3, 0.3, num, 11, GRAY3, True)
        add_text(s, left, 2.85, 2.3, 0.25, loc, 10, GRAY3)
        add_text(s, left, 3.15, 2.3, 0.35, title, 13, BLACK, True)
        add_text(s, left, 3.55, 2.3, 0.7, desc, 11, GRAY2)
    
    add_text(s, 0.9, 5.0, 5, 0.35, "Core 18-week sprint", 14, BLACK, True)
    add_text(s, 6, 5.0, 6, 0.35, "Deep build alongside FOUNTAIN", 13, GRAY2)
    
    # ═══════════ SLIDE 6: PHASE DETAILS ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "PROGRAM DETAILS", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.5, "What happens at each phase", 32, BLACK, True)
    
    details = [
        ("Phase 1 · SF", "The Workshop & Alignment", "Turning ideas into a runnable intelligence distribution path. Solving for Product–Model Fit, Product–Market Fit, Product–MP Fit."),
        ("Phase 2 · SZ", "The Immersion & Reality", "Learn how the physical world runs—materials, processes, costs, yields, certifications. Work from the FOUNTAIN Physical AI House."),
        ("Phase 3 · SZ & SF", "The Build & Merge", "Work with FOUNTAIN community and Advisors. Combine AI and Physical architectures into one system."),
        ("Phase 4-5 · SZ → SF", "EVT Output & Launch", "Engineering Validation Test units. Mechanical, electronics, AI integration. User trials, KOL testing, brand story.")
    ]
    for i, (phase, title, desc) in enumerate(details):
        row, col = i // 2, i % 2
        add_text(s, 0.9 + col * 6, 2.0 + row * 1.7, 5.5, 0.25, phase, 11, GRAY3)
        add_text(s, 0.9 + col * 6, 2.3 + row * 1.7, 5.5, 0.35, title, 15, BLACK, True)
        add_text(s, 0.9 + col * 6, 2.7 + row * 1.7, 5.5, 1, desc, 13, GRAY)
    
    # ═══════════ SLIDE 7: LAUNCH & SCALE ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "MONTH 4-6", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "Launch Day is Graduation Day.", 40, BLACK, True)
    add_text(s, 0.9, 1.9, 11, 0.4, "We don't run a demo day. We ship the Physical AI.", 18, GRAY2)
    add_text(s, 0.9, 2.5, 11, 0.35, "From there, three scale engines activate:", 15, GRAY3)
    
    engines = [
        ("1", "Mass Production", "First production run (≈3K units) from EVT to real distribution."),
        ("2", "Growth System", "Hardware: crowdfunding, DTC, e-commerce, distributors.\nSoftware: Product Hunt, dev communities, social."),
        ("3", "Iteration Loop", "Real-world usage drives data, data drives model evolution, models drive next product cycle.")
    ]
    for i, (num, title, desc) in enumerate(engines):
        left = 0.9 + i * 4
        add_text(s, left, 3.2, 0.5, 0.5, num, 24, BLACK, True)
        add_text(s, left, 3.7, 3.5, 0.35, title, 14, BLACK, True)
        add_text(s, left, 4.1, 3.5, 1, desc, 13, GRAY)
    
    add_text(s, 0, 5.5, 13.333, 0.4, "Silicon Valley thinking × Shenzhen execution", 14, GRAY3, align=PP_ALIGN.CENTER)
    
    # ═══════════ SLIDE 8: PARTNERSHIP ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "HOW WE PARTNER", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "Two paths to work with FOUNTAIN", 40, BLACK, True)
    
    # Path 1
    add_text(s, 0.9, 2.4, 5.5, 0.25, "PATH 1", 11, GRAY3, True)
    add_text(s, 0.9, 2.7, 5.5, 0.4, "The Full Stack", 18, BLACK, True)
    add_text(s, 0.9, 3.15, 5.5, 0.6, "$500K", 48, BLACK, True)
    add_text(s, 0.9, 3.8, 5.5, 0.35, "for 10% equity", 16, GRAY)
    add_text(s, 0.9, 4.3, 5.5, 0.7, "5% → FOUNTAIN (Execution System)\n5% → Capital Partner", 13, GRAY2)
    
    # Path 2
    add_text(s, 7, 2.4, 5.5, 0.25, "PATH 2", 11, GRAY3, True)
    add_text(s, 7, 2.7, 5.5, 0.4, "The Execution Engine", 18, BLACK, True)
    add_text(s, 7, 3.15, 5.5, 0.6, "5%", 48, BLACK, True)
    add_text(s, 7, 3.8, 5.5, 0.35, "equity (Program Only)", 16, GRAY)
    add_text(s, 7, 4.3, 5.5, 0.7, "5% → FOUNTAIN (Execution System)\nFor funded teams who need execution.", 13, GRAY2)
    
    add_text(s, 0, 5.5, 13.333, 0.4, "We work with a small number of teams each cycle.", 14, GRAY3, align=PP_ALIGN.CENTER)
    
    # ═══════════ SLIDE 9: BATCH F1 ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "BUILDERS WITH US", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "FOUNTAIN Batch F1", 40, BLACK, True)
    add_text(s, 0.9, 1.85, 11, 0.35, "First Batch of builders is moving into the world.", 15, GRAY3)
    
    batch = [
        ("🔖", "Mark", "Los Angeles", "Eason Tang", "AI Bookmark for reading and knowledge", "mark.engineering", "Q2 2026"),
        ("💍", "Sparkring", "Shenzhen", "Charley Tang", "AI Ring for voice workflow", "sparkring.ai", "Q2 2026"),
        ("🩺", "Cortexa", "Palo Alto", "Jack Lee", "AI Medical Device for clinical documentation", "cortexanote.com", "Q1 2026")
    ]
    for i, (emoji, name, loc, founder, desc, url, launch) in enumerate(batch):
        left = 0.9 + i * 4
        add_text(s, left, 2.6, 0.5, 0.4, emoji, 24)
        add_text(s, left + 2.5, 2.6, 1, 0.3, launch, 10, GRAY3)
        add_text(s, left, 3.05, 3.5, 0.4, name, 20, BLACK, True)
        add_text(s, left, 3.45, 3.5, 0.25, loc, 12, GRAY3)
        add_text(s, left, 3.75, 3.5, 0.3, f"Founder: {founder}", 13, GRAY2)
        add_text(s, left, 4.1, 3.5, 0.5, desc, 14, GRAY)
        add_text(s, left, 4.6, 3.5, 0.25, url, 12, GRAY3)
    
    # ═══════════ SLIDE 10: TEAM ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "WHO WE ARE", 11, GRAY3, True)
    add_text(s, 0.9, 1.2, 11, 0.6, "Built by builders.", 40, BLACK, True)
    add_text(s, 0.9, 1.85, 11, 0.5, "For builders building in the real world.\nOperating between Silicon Valley and Shenzhen.", 16, GRAY2)
    
    # Founder
    add_text(s, 0.9, 2.8, 5.5, 0.45, "🌲 Tree Wan", 22, BLACK, True)
    add_text(s, 0.9, 3.3, 5.5, 0.25, "Cofounder", 13, GRAY3)
    add_text(s, 0.9, 3.7, 5.5, 1.2, "Serial software & robotics entrepreneur.\nFounder of Speedfox / See / Zeustech.\nDecade of software & hardware delivery to millions of users.", 14, GRAY)
    add_text(s, 0.9, 5.0, 5.5, 0.3, "Backed by Sequoia · IDG · BAI · Tencent · Alibaba", 12, GRAY3)
    
    # Advisory
    add_text(s, 7, 2.8, 5.5, 0.35, "🤖 ADVISORY BOARD", 11, GRAY3, True)
    advisors = "DG Robotics\nDobot\nRobosen\nNarwal Robot\nGoeroptics\nLiepin\nWatermark Camera\nNeabot"
    add_text(s, 7, 3.3, 5.5, 2.5, advisors, 14, GRAY)
    
    # ═══════════ SLIDE 11: CTA ═══════════
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, 0.9, 0.7, 3, 0.3, "JOIN THE BUILD", 11, GRAY3, True)
    add_text(s, 0.9, 1.8, 11, 1, "If you're building Physical AI,\nwe should talk.", 48, BLACK, True)
    add_text(s, 0.9, 3.5, 11, 0.5, "Thinking about real-world deployment?\nLet's have a conversation.", 18, GRAY2)
    
    add_text(s, 0.9, 4.5, 5, 0.25, "30-minute conversation", 14, GRAY3, True)
    add_text(s, 0.9, 4.8, 10, 0.35, "Team · AI · Product · Users · Bay Area × Shenzhen", 16, GRAY)
    
    add_text(s, 0.9, 5.5, 5, 0.4, "apply@fountainbuild.ai", 24, BLACK, True)
    add_text(s, 0.9, 6.0, 5, 0.3, "fountainbuild.ai", 16, GRAY3)
    
    add_text(s, 0.9, 6.8, 5, 0.25, "🔹 FOUNTAIN", 12, GRAY3, True)
    
    # Save
    prs.save(os.path.join(os.path.dirname(__file__), 'FOUNTAIN_Pitchdeck_Brand.pptx'))
    print("PPT generated")

if __name__ == "__main__":
    create_deck()
