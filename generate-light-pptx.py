#!/usr/bin/env python3
"""
Generate FOUNTAIN Pitch Deck - Light/Minimal Version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors - Minimal palette
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(15, 23, 42)  # #0f172a
GRAY = RGBColor(71, 85, 105)  # #475569
LIGHT_GRAY = RGBColor(148, 163, 184)  # #94a3b8
LINE = RGBColor(226, 232, 240)  # #e2e8f0

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                      Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return shape

def create_pitchdeck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "FOUNTAIN", 
                 font_size=13, font_color=BLACK, bold=True)
    
    add_text_box(slide, 0.9, 2.5, 11, 1.5, "Physical AI\nAccelerator", 
                 font_size=72, font_color=BLACK, bold=True)
    
    add_text_box(slide, 0.9, 4.8, 10, 0.6, "From AI demo to shipped product. 6 months.", 
                 font_size=24, font_color=GRAY)
    
    add_text_box(slide, 0.9, 5.6, 10, 0.4, "Bay Area × Shenzhen", 
                 font_size=14, font_color=LIGHT_GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "THE PROBLEM", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 11, 1.2, 
                 "AI has scaled rapidly in the digital world.\nThe physical world is just beginning.", 
                 font_size=44, font_color=BLACK, bold=True)
    
    problems = [
        ("Model ≠ Product", "Building AI models or apps is not the same as shipping physical systems."),
        ("Demo → Production Gap", "The gap between demo and mass production is vast."),
        ("East-West Disconnect", "Silicon Valley has ideas. Shenzhen has factories. Few bridge both.")
    ]
    
    for i, (title, desc) in enumerate(problems):
        left = 0.9 + i * 4
        add_text_box(slide, left, 3.8, 3.5, 0.4, title, font_size=16, font_color=BLACK, bold=True)
        add_text_box(slide, left, 4.35, 3.5, 1, desc, font_size=15, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 3: SOLUTION
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "OUR SOLUTION", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 1, 
                 "We don't just advise.\nWe build alongside founders.", 
                 font_size=44, font_color=BLACK, bold=True)
    
    add_text_box(slide, 0.9, 2.8, 10, 0.5, 
                 "A 6-month sprint from San Francisco to Shenzhen. Launch Day is Graduation Day.", 
                 font_size=18, font_color=GRAY)
    
    solutions = [
        ("Hands-On Building", "Hardware design, firmware, supply chain, factory relationships."),
        ("Global Network", "Deep connections in Bay Area and Shenzhen."),
        ("Ship, Not Demo", "Real products to real customers."),
        ("Capital + Execution", "Funding paired with operational support.")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        row = i // 2
        col = i % 2
        left = 0.9 + col * 6
        top = 3.8 + row * 1.3
        add_text_box(slide, left, top, 5.5, 0.35, title, font_size=16, font_color=BLACK, bold=True)
        add_text_box(slide, left, top + 0.4, 5.5, 0.6, desc, font_size=15, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 4: WHY NOW
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "WHY NOW", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 0.7, 
                 "The Physical AI moment has arrived.", 
                 font_size=44, font_color=BLACK, bold=True)
    
    stats = [
        ("10×", "Cost Reduction", "AI compute costs dropped 10× in 3 years"),
        ("$2T", "Market by 2030", "Physical AI market projected (McKinsey)"),
        ("Now", "Timing", "First movers will define next decade")
    ]
    
    for i, (num, unit, desc) in enumerate(stats):
        left = 0.9 + i * 4
        add_text_box(slide, left, 3, 3.5, 0.9, num, font_size=56, font_color=BLACK, bold=True)
        add_text_box(slide, left, 3.9, 3.5, 0.35, unit, font_size=14, font_color=GRAY, bold=True)
        add_text_box(slide, left, 4.35, 3.5, 0.6, desc, font_size=14, font_color=LIGHT_GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 5: PROGRAM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "THE PROGRAM", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 0.7, 
                 "6 months. Silicon Valley ↔ Shenzhen.", 
                 font_size=44, font_color=BLACK, bold=True)
    
    steps = [
        ("01", "San Francisco", "PMF Alignment", "Validate product-market fit"),
        ("02", "Shenzhen", "Ecosystem Immersion", "Factory tours, supplier intros"),
        ("03", "Shenzhen", "Hardware Build", "ID/MD, PCB, firmware"),
        ("04", "Shenzhen", "EVT Prototypes", "Engineering validation"),
        ("05", "San Francisco", "Launch Day", "Ship to customers")
    ]
    
    for i, (num, loc, title, desc) in enumerate(steps):
        left = 0.6 + i * 2.5
        add_text_box(slide, left, 2.8, 2.3, 0.5, num, font_size=32, font_color=BLACK, bold=True)
        add_text_box(slide, left, 3.4, 2.3, 0.3, loc, font_size=11, font_color=LIGHT_GRAY, bold=True)
        add_text_box(slide, left, 3.75, 2.3, 0.35, title, font_size=15, font_color=BLACK, bold=True)
        add_text_box(slide, left, 4.15, 2.3, 0.5, desc, font_size=13, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 6: PORTFOLIO
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "BATCH F1", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 0.7, 
                 "First builders moving into the world", 
                 font_size=44, font_color=BLACK, bold=True)
    
    portfolio = [
        ("🔖", "Mark", "AI Bookmark for reading and knowledge", "Eason Tang · Los Angeles", "Q2 2026"),
        ("💍", "Sparkring", "AI Ring for voice-powered workflow", "Charley Tang · Shenzhen", "Q2 2026"),
        ("🩺", "Cortexa", "AI Medical Device for clinical docs", "Jack Lee · Palo Alto", "Q1 2026")
    ]
    
    for i, (icon, name, tagline, meta, launch) in enumerate(portfolio):
        left = 0.9 + i * 4
        add_text_box(slide, left, 2.8, 0.5, 0.4, icon, font_size=24)
        add_text_box(slide, left + 2.5, 2.8, 1, 0.3, launch, font_size=11, font_color=LIGHT_GRAY)
        add_text_box(slide, left, 3.3, 3.5, 0.4, name, font_size=20, font_color=BLACK, bold=True)
        add_text_box(slide, left, 3.75, 3.5, 0.5, tagline, font_size=14, font_color=GRAY)
        add_text_box(slide, left, 4.35, 3.5, 0.3, meta, font_size=12, font_color=LIGHT_GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 7: TEAM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "TEAM", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 0.7, 
                 "Operators who've shipped at scale", 
                 font_size=44, font_color=BLACK, bold=True)
    
    # Founder
    add_text_box(slide, 0.9, 2.6, 5.5, 0.45, "Tree Wan", font_size=24, font_color=BLACK, bold=True)
    add_text_box(slide, 0.9, 3.1, 5.5, 0.3, "Co-Founder", font_size=14, font_color=LIGHT_GRAY)
    add_text_box(slide, 0.9, 3.6, 5.5, 1.4, 
                 "Founder of Speedfox, See, Zeustech. A decade of shipping software and hardware products to millions of users. Former Tencent PM with 13 international patents.\n\nForbes 30 Under 30 Asia · Retail & E-commerce, 2017", 
                 font_size=14, font_color=GRAY)
    add_text_box(slide, 0.9, 5.4, 5.5, 0.3, "Sequoia · IDG · BAI · Tencent · Alibaba", 
                 font_size=12, font_color=LIGHT_GRAY)
    
    # Advisory
    add_text_box(slide, 7.5, 2.6, 5, 0.3, "ADVISORY NETWORK", font_size=11, font_color=LIGHT_GRAY, bold=True)
    advisors = "DG Robotics\nDobot\nRobosen\nNarwal Robot\nGoeroptics\nNeabot"
    add_text_box(slide, 7.5, 3.1, 5, 2.5, advisors, font_size=15, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 8: INVESTMENT
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 0.8, 3, 0.4, "INVESTMENT", 
                 font_size=11, font_color=LIGHT_GRAY, bold=True)
    
    add_text_box(slide, 0.9, 1.4, 10, 0.7, 
                 "Two ways to work with Fountain", 
                 font_size=44, font_color=BLACK, bold=True)
    
    # Option 1
    add_text_box(slide, 0.9, 2.5, 5.5, 0.3, "FULL PARTNERSHIP", font_size=11, font_color=LIGHT_GRAY, bold=True)
    add_text_box(slide, 0.9, 2.9, 5.5, 0.6, "$500K", font_size=48, font_color=BLACK, bold=True)
    add_text_box(slide, 0.9, 3.55, 5.5, 0.35, "for 10% equity", font_size=18, font_color=GRAY)
    benefits1 = "· 6-month program (SF ↔ Shenzhen)\n· Hardware development support\n· Supply chain & factory access\n· Shenzhen office & local team\n· Investor network introductions\n· Launch day execution support"
    add_text_box(slide, 0.9, 4.1, 5.5, 2, benefits1, font_size=13, font_color=GRAY)
    
    # Option 2
    add_text_box(slide, 7, 2.5, 5.5, 0.3, "EXECUTION ONLY", font_size=11, font_color=LIGHT_GRAY, bold=True)
    add_text_box(slide, 7, 2.9, 5.5, 0.6, "5%", font_size=48, font_color=BLACK, bold=True)
    add_text_box(slide, 7, 3.55, 5.5, 0.35, "equity (no cash)", font_size=18, font_color=GRAY)
    benefits2 = "· 6-month program (SF ↔ Shenzhen)\n· Hardware development support\n· Supply chain & factory access\n· Shenzhen office & local team\n· For funded teams who need execution"
    add_text_box(slide, 7, 4.1, 5.5, 2, benefits2, font_size=13, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 9: CTA
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    
    add_text_box(slide, 0.9, 2.2, 11, 1, "Building Physical AI?\nLet's talk.", 
                 font_size=56, font_color=BLACK, bold=True)
    
    add_text_box(slide, 0.9, 4, 11, 0.4, "A 30-minute conversation. Team · AI · Product · Users.", 
                 font_size=18, font_color=GRAY)
    
    add_text_box(slide, 0.9, 4.8, 11, 0.45, "apply@fountainbuild.ai", 
                 font_size=28, font_color=BLACK, bold=True)
    
    add_text_box(slide, 0.9, 5.4, 11, 0.35, "fountainbuild.ai", 
                 font_size=16, font_color=LIGHT_GRAY)
    
    add_text_box(slide, 0.9, 6.5, 11, 0.3, "FOUNTAIN · Physical AI Accelerator · Bay Area × Shenzhen", 
                 font_size=13, font_color=LIGHT_GRAY)
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), 'FOUNTAIN_Pitchdeck_Light.pptx')
    prs.save(output_path)
    print(f"PowerPoint generated: {output_path}")

if __name__ == "__main__":
    create_pitchdeck()
