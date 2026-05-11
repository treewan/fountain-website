#!/usr/bin/env python3
"""
Generate FOUNTAIN Pitch Deck as PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG_DARK = RgbColor(10, 15, 26)  # #0a0f1a
BLUE = RgbColor(37, 99, 235)  # #2563eb
BLUE_L = RgbColor(96, 165, 250)  # #60a5fa
CYAN = RgbColor(6, 182, 212)  # #06b6d4
GREEN = RgbColor(16, 185, 129)  # #10b981
ORANGE = RgbColor(245, 158, 11)  # #f59e0b
WHITE = RgbColor(248, 250, 252)  # #f8fafc
GRAY = RgbColor(148, 163, 184)  # #94a3b8
DARK_GRAY = RgbColor(100, 116, 139)  # #64748b

def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to slide"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                      Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return shape

def create_pitchdeck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 1: COVER
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Logo
    add_text_box(slide, 0, 0.8, 13.333, 0.5, "🔹 FOUNTAIN", 
                 font_size=14, font_color=BLUE_L, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Main title
    add_text_box(slide, 0, 2.2, 13.333, 1.5, "Physical AI\nAccelerator", 
                 font_size=72, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Tagline
    add_text_box(slide, 0, 4.3, 13.333, 0.6, "From AI demo to shipped product. 6 months.", 
                 font_size=24, font_color=GRAY, alignment=PP_ALIGN.CENTER)
    
    # Meta
    add_text_box(slide, 0, 5.3, 13.333, 0.5, "Bay Area  ×  Shenzhen", 
                 font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "THE PROBLEM", 
                 font_size=11, font_color=ORANGE, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 11, 1.5, 
                 "AI has scaled rapidly in the digital world.\nThe physical world is just beginning.", 
                 font_size=42, font_color=WHITE, bold=True)
    
    # Problem cards
    problems = [
        ("🧠", "Model ≠ Product", "Building AI models or apps is not the same as shipping physical systems."),
        ("🏭", "Demo → Production Gap", "The gap between demo and mass production is vast."),
        ("🌏", "East-West Disconnect", "Silicon Valley has ideas. Shenzhen has factories. Few bridge both.")
    ]
    
    for i, (icon, title, desc) in enumerate(problems):
        left = 0.8 + i * 4
        add_text_box(slide, left, 3.5, 0.5, 0.5, icon, font_size=28)
        add_text_box(slide, left, 4.1, 3.5, 0.4, title, font_size=18, font_color=WHITE, bold=True)
        add_text_box(slide, left, 4.6, 3.5, 1, desc, font_size=14, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 3: SOLUTION
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "OUR SOLUTION", 
                 font_size=11, font_color=GREEN, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 1.2, 
                 "We don't just advise.\nWe build alongside founders.", 
                 font_size=42, font_color=WHITE, bold=True)
    
    add_text_box(slide, 0.8, 2.5, 10, 0.5, 
                 "A 6-month sprint from San Francisco to Shenzhen. Launch Day is Graduation Day.", 
                 font_size=18, font_color=GRAY)
    
    solutions = [
        ("🛠️ Hands-On Building", "Hardware design, firmware, supply chain, factory relationships."),
        ("🌐 Global Network", "Deep connections in Bay Area and Shenzhen."),
        ("📦 Ship, Not Demo", "Real products to real customers. Launch is graduation."),
        ("💰 Capital + Execution", "Funding paired with operational support.")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        row = i // 2
        col = i % 2
        left = 0.8 + col * 6
        top = 3.4 + row * 1.5
        add_text_box(slide, left, top, 5.5, 0.4, title, font_size=18, font_color=WHITE, bold=True)
        add_text_box(slide, left, top + 0.45, 5.5, 0.8, desc, font_size=14, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 4: WHY NOW
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "WHY NOW", 
                 font_size=11, font_color=CYAN, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 0.8, 
                 "The Physical AI moment has arrived.", 
                 font_size=42, font_color=WHITE, bold=True)
    
    stats = [
        ("10×", "Cost Reduction", "AI compute costs dropped 10× in 3 years"),
        ("$2T", "Market by 2030", "Physical AI market projected (McKinsey)"),
        ("Now", "Timing", "First movers will define next decade")
    ]
    
    for i, (num, unit, desc) in enumerate(stats):
        left = 0.8 + i * 4
        add_text_box(slide, left, 2.8, 3.5, 1, num, font_size=56, font_color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, 4.0, 3.5, 0.4, unit, font_size=14, font_color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, 4.5, 3.5, 0.8, desc, font_size=13, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 5: PROGRAM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "THE PROGRAM", 
                 font_size=11, font_color=BLUE_L, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 0.8, 
                 "6 months. Silicon Valley ↔ Shenzhen.", 
                 font_size=42, font_color=WHITE, bold=True)
    
    steps = [
        ("1", "San Francisco", "PMF Alignment", "Validate product-market fit"),
        ("2", "Shenzhen", "Ecosystem Immersion", "Factory tours, supplier intros"),
        ("3", "Shenzhen", "Hardware Build", "ID/MD design, PCB, firmware"),
        ("4", "Shenzhen", "EVT Prototypes", "Engineering validation"),
        ("5", "San Francisco", "Launch Day", "Ship to customers")
    ]
    
    for i, (num, loc, title, desc) in enumerate(steps):
        left = 0.6 + i * 2.5
        add_text_box(slide, left, 2.8, 0.6, 0.6, num, font_size=24, font_color=BLUE_L, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, 3.6, 2.3, 0.3, loc, font_size=10, font_color=BLUE_L, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, 3.95, 2.3, 0.4, title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, 4.4, 2.3, 0.6, desc, font_size=11, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 6: PORTFOLIO
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "BATCH F1", 
                 font_size=11, font_color=BLUE_L, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 0.8, 
                 "First builders moving into the world", 
                 font_size=42, font_color=WHITE, bold=True)
    
    portfolio = [
        ("🔖", "Mark", "AI Bookmark for reading and knowledge", "Eason Tang · Los Angeles", "Q2 2026"),
        ("💍", "Sparkring", "AI Ring for voice-powered workflow", "Charley Tang · Shenzhen", "Q2 2026"),
        ("🩺", "Cortexa", "AI Medical Device for clinical docs", "Jack Lee · Palo Alto", "Q1 2026")
    ]
    
    for i, (icon, name, tagline, meta, launch) in enumerate(portfolio):
        left = 0.8 + i * 4
        add_text_box(slide, left, 2.5, 0.6, 0.5, icon, font_size=32)
        add_text_box(slide, left + 2.5, 2.5, 1, 0.4, launch, font_size=10, font_color=GREEN, bold=True)
        add_text_box(slide, left, 3.1, 3.5, 0.5, name, font_size=22, font_color=WHITE, bold=True)
        add_text_box(slide, left, 3.6, 3.5, 0.6, tagline, font_size=14, font_color=GRAY)
        add_text_box(slide, left, 4.3, 3.5, 0.4, meta, font_size=12, font_color=DARK_GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 7: TEAM
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "TEAM", 
                 font_size=11, font_color=BLUE_L, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 0.8, 
                 "Operators who've shipped hardware at scale", 
                 font_size=42, font_color=WHITE, bold=True)
    
    # Founder info
    add_text_box(slide, 0.8, 2.4, 6, 0.5, "Tree Wan", font_size=28, font_color=WHITE, bold=True)
    add_text_box(slide, 0.8, 2.95, 6, 0.4, "Co-Founder", font_size=14, font_color=BLUE_L, bold=True)
    add_text_box(slide, 0.8, 3.5, 5.5, 1.2, 
                 "Founder of Speedfox, See, Zeustech. A decade of shipping software and hardware products to millions of users. Former Tencent PM with 13 international patents.\n\nForbes 30 Under 30 Asia (Retail & E-commerce, 2017)", 
                 font_size=13, font_color=GRAY)
    
    # Backers
    add_text_box(slide, 0.8, 5.2, 6, 0.4, "Backed by: Sequoia · IDG · BAI · Tencent · Alibaba", 
                 font_size=12, font_color=DARK_GRAY)
    
    # Advisory
    add_text_box(slide, 7.5, 2.4, 5, 0.4, "ADVISORY NETWORK", font_size=11, font_color=DARK_GRAY, bold=True)
    advisors = "DG Robotics\nDobot\nRobosen\nNarwal Robot\nGoeroptics\nNeabot"
    add_text_box(slide, 7.5, 2.9, 5, 2.5, advisors, font_size=16, font_color=WHITE)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 8: INVESTMENT
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0.8, 0.6, 3, 0.4, "INVESTMENT", 
                 font_size=11, font_color=GREEN, bold=True)
    
    add_text_box(slide, 0.8, 1.1, 10, 0.8, 
                 "Two ways to work with Fountain", 
                 font_size=42, font_color=WHITE, bold=True)
    
    # Option 1: Full Partnership
    add_text_box(slide, 0.8, 2.3, 5.5, 0.4, "FULL PARTNERSHIP (RECOMMENDED)", font_size=12, font_color=GRAY, bold=True)
    add_text_box(slide, 0.8, 2.8, 5.5, 0.7, "$500K", font_size=48, font_color=BLUE_L, bold=True)
    add_text_box(slide, 0.8, 3.5, 5.5, 0.4, "for 10% equity", font_size=18, font_color=GRAY)
    benefits1 = "✓ 6-month program (SF ↔ Shenzhen)\n✓ Hardware development support\n✓ Supply chain & factory access\n✓ Shenzhen office & local team\n✓ Investor network introductions\n✓ Launch day execution support"
    add_text_box(slide, 0.8, 4.1, 5.5, 2, benefits1, font_size=13, font_color=GRAY)
    
    # Option 2: Execution Only
    add_text_box(slide, 7, 2.3, 5.5, 0.4, "EXECUTION ONLY", font_size=12, font_color=GRAY, bold=True)
    add_text_box(slide, 7, 2.8, 5.5, 0.7, "5%", font_size=48, font_color=WHITE, bold=True)
    add_text_box(slide, 7, 3.5, 5.5, 0.4, "equity (no cash)", font_size=18, font_color=GRAY)
    benefits2 = "✓ 6-month program (SF ↔ Shenzhen)\n✓ Hardware development support\n✓ Supply chain & factory access\n✓ Shenzhen office & local team\n✓ For funded teams who need execution"
    add_text_box(slide, 7, 4.1, 5.5, 2, benefits2, font_size=13, font_color=GRAY)
    
    # ═══════════════════════════════════════════════════════════
    # SLIDE 9: CTA
    # ═══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    add_text_box(slide, 0, 2, 13.333, 1, "Building Physical AI?\nLet's talk.", 
                 font_size=48, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0, 3.8, 13.333, 0.5, "A 30-minute conversation. Team · AI · Product · Users.", 
                 font_size=18, font_color=GRAY, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0, 4.8, 13.333, 0.5, "apply@fountainbuild.ai", 
                 font_size=24, font_color=BLUE_L, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0, 5.4, 13.333, 0.4, "fountainbuild.ai", 
                 font_size=16, font_color=GRAY, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, 0, 6.5, 13.333, 0.4, "🔹 FOUNTAIN · Physical AI Accelerator · Bay Area × Shenzhen", 
                 font_size=13, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), 'FOUNTAIN_Pitchdeck.pptx')
    prs.save(output_path)
    print(f"PowerPoint generated: {output_path}")

if __name__ == "__main__":
    create_pitchdeck()
